"""Fix Dongguan showroom products: correct categories, image mapping, and orientation."""
import sys, io, os, re, zipfile
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = r"C:\Users\eukri\OneDrive\Documents\Claude Code\Credentials Claude Code\ai-agents-go-4c81b70995db.json"

from pptx import Presentation
from pptx.util import Emu
from xml.etree import ElementTree as ET
from google.cloud.firestore_v1 import Client
from datetime import datetime, timezone

db = Client(project="ai-agents-go", database="areda-product-catalogs")
now = datetime.now(timezone.utc).isoformat()

EXCHANGE_RATE = 7.2  # RMB per USD
RETAIL_MARKUP = 1.2  # 20% international markup

def compute_retail_usd(rmb: float) -> float:
    """Auto-convert RetailPriceRMB to RetailPriceUSD: RMB / exchangeRate * 1.2"""
    if rmb > 0:
        return round(rmb / EXCHANGE_RATE * RETAIL_MARKUP, 2)
    return 0
pptx_path = r"C:\Users\eukri\OneDrive\Documents\Documents GO\WeChat OneDrive\WeChat Grant visconti\2026-02-10 卓越卡萨东莞展厅现货汇总-2025-1110.pptx"
prs = Presentation(pptx_path)
base_url = "https://storage.googleapis.com/areda-product-images/dongguan-showroom"

# Build slide -> images with positions
slide_images = {}
with zipfile.ZipFile(pptx_path) as z:
    for slide_num in range(1, len(prs.slides) + 1):
        rels_file = f"ppt/slides/_rels/slide{slide_num}.xml.rels"
        if rels_file not in z.namelist():
            continue
        tree = ET.parse(z.open(rels_file))
        ns = "{http://schemas.openxmlformats.org/package/2006/relationships}"
        rid_to_img = {}
        for rel in tree.getroot().findall(f"{ns}Relationship"):
            target = rel.get("Target", "")
            rid = rel.get("Id", "")
            if "../media/" in target:
                rid_to_img[rid] = target.split("/")[-1]

        # Get image shapes with positions from slide XML
        slide_xml = f"ppt/slides/slide{slide_num}.xml"
        if slide_xml not in z.namelist():
            continue
        stree = ET.parse(z.open(slide_xml))
        sroot = stree.getroot()

        imgs_with_pos = []
        for elem in sroot.iter():
            tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
            if tag == "pic":
                # Find blipFill -> blip -> r:embed
                blip = None
                x, y, w, h = 0, 0, 0, 0
                for child in elem.iter():
                    ctag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                    if ctag == "blip":
                        for k, v in child.attrib.items():
                            if k.endswith("}embed"):
                                blip = v
                    if ctag == "off":
                        x = int(child.get("x", 0))
                        y = int(child.get("y", 0))
                    if ctag == "ext":
                        w = int(child.get("cx", 0))
                        h = int(child.get("cy", 0))

                if blip and blip in rid_to_img:
                    img_name = rid_to_img[blip]
                    # Filter out tiny images (logos, icons) - less than ~100x100px
                    if w > 500000 and h > 500000:  # EMU units, ~0.5cm threshold
                        imgs_with_pos.append({
                            "name": img_name,
                            "x": x, "y": y, "w": w, "h": h,
                            "area": w * h
                        })

        # Sort by position (top-left to bottom-right, prioritize larger images)
        imgs_with_pos.sort(key=lambda i: (i["y"], i["x"]))
        slide_images[slide_num] = imgs_with_pos


def categorize(name_cn, text_block=""):
    """Better categorization based on Chinese product names."""
    n = name_cn
    if "单人沙发" in n or "单人" in n:
        return "lounge-chairs"
    if "沙发" in n:
        return "sofas"
    if "餐椅" in n:
        return "dining-chairs"
    if "书椅" in n or "办公椅" in n:
        return "lounge-chairs"
    if "吧台椅" in n or "吧椅" in n:
        return "dining-chairs"
    if "椅子" in n or "椅" in n:
        return "dining-chairs"
    if "床" in n and "头柜" not in n and "尾榻" not in n:
        return "beds"
    if "床头柜" in n:
        return "storage"
    if "茶几" in n or "边几" in n or "角几" in n:
        return "tables"
    if "餐桌" in n or "书桌" in n or "茶桌" in n:
        return "tables"
    if "柜" in n or "展架" in n:
        return "storage"
    if "墩" in n or "脚凳" in n:
        return "sofas"
    if "榻" in n:
        return "sofas"
    if "长椅" in n or "凳" in n:
        return "sofas"
    # Fallback: check the full text block for clues
    if "床垫" in text_block or "床" in text_block:
        return "beds"
    if "沙发" in text_block:
        return "sofas"
    return "lounge-chairs"  # Default for unrecognized single items


# Delete existing Dongguan products
print("Deleting existing Dongguan products...")
from google.cloud.firestore_v1.base_query import FieldFilter
dongguan_docs = list(db.collection("products").where(
    filter=FieldFilter("collection", "==", "Dongguan Showroom")
).stream())
for doc in dongguan_docs:
    db.collection("products").document(doc.id).delete()
print(f"Deleted {len(dongguan_docs)} existing Dongguan products")

# Re-extract products with correct image mapping
products = []
for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text += shape.text_frame.text + "\n"

    # Split into product blocks by price marker
    blocks = re.split(r"(?=(?:指导价|一口价)\s*[￥¥])", text)

    slide_imgs = slide_images.get(slide_num, [])
    product_idx = 0

    for block in blocks:
        name_m = re.search(r"名称[：:]\s*(.+?)(?:\n|$)", block)
        if not name_m:
            continue
        raw_name = name_m.group(1).strip()
        if not raw_name or raw_name == "维斯康缇":
            continue

        name_cn = re.sub(r"\s*(已售|编码.*)", "", raw_name).strip()

        code_m = re.search(r"编码[：:]\s*([A-Za-z0-9\-]+)", block)
        code = code_m.group(1).strip() if code_m else f"DG-S{slide_num:02d}-{product_idx+1}"

        price_m = re.search(r"[￥¥]\s*([\d,.]+)", block)
        price_rmb = float(price_m.group(1).replace(",", "")) if price_m else 0

        # Check for multiplier (e.g., "￥8576*2个")
        mult_m = re.search(r"[￥¥]\s*[\d,.]+\s*\*\s*(\d+)", block)
        qty_note = f" (x{mult_m.group(1)})" if mult_m else ""

        dims_m = re.search(r"尺寸[：:]\s*(.+?)(?:\n|$)", block)
        dims = dims_m.group(1).strip() if dims_m else ""

        mat_m = re.search(r"材质[：:]\s*(.+?)(?:\n\n|\n(?:指导|一口|名称)|$)", block, re.DOTALL)
        material = mat_m.group(1).strip().replace("\n", " ").replace("          ", " ")[:200] if mat_m else ""

        category = categorize(name_cn, block)

        # Assign image: try to match product to its closest image by position
        img_url = ""
        if slide_imgs:
            if len(slide_imgs) == 1:
                # Single image on slide - all products share it
                img_url = f"{base_url}/{slide_imgs[0]['name']}"
            elif product_idx < len(slide_imgs):
                # Multiple images - assign by order (top to bottom)
                img_url = f"{base_url}/{slide_imgs[product_idx]['name']}"
            else:
                # More products than images - use the largest image
                biggest = max(slide_imgs, key=lambda x: x["area"])
                img_url = f"{base_url}/{biggest['name']}"

        # Check if sold
        is_sold = "已售" in raw_name or "已售" in block

        products.append({
            "productCode": code,
            "brand": "VISCONTI",
            "style": "Dongguan Showroom",
            "styleCn": "",
            "name": f"{name_cn}{qty_note} (Showroom)",
            "nameCn": name_cn,
            "model": code,
            "dimensions": dims,
            "material": material,
            "materialCn": material,
            "features": "已售 (Sold)" if is_sold else "现货 (In Stock)",
            "featuresCn": "",
            "category": category,
            "imageUrl": img_url,
            "galleryUrls": [],
            "volumeCbm": 0,
            "priceExwUsd": 0,
            "priceFobUsd": 0,
            "priceRmb": price_rmb,
            "retailPriceRmb": price_rmb,
            "retailPriceUsd": compute_retail_usd(price_rmb),
            "exchangeRate": EXCHANGE_RATE,
            "currency": "RMB",
            "moq": 1,
            "leadTimeDays": 0,
            "collection": "Dongguan Showroom",
            "tags": ["showroom", "sold" if is_sold else "in-stock"],
            "visible": True,
            "createdAt": now,
            "updatedAt": now,
        })
        product_idx += 1

print(f"\nRe-seeding {len(products)} Dongguan products...")

# Category summary
from collections import Counter
cats = Counter(p["category"] for p in products)
print(f"Categories: {dict(cats)}")

with_img = sum(1 for p in products if p["imageUrl"])
print(f"With images: {with_img}/{len(products)}")

sold = sum(1 for p in products if "sold" in p["tags"])
print(f"Sold items: {sold}")

# Batch write
batch = db.batch()
for i, p in enumerate(products):
    doc_id = p["productCode"].lower().replace(" ", "-")
    ref = db.collection("products").document(doc_id)
    batch.set(ref, p)
    if (i + 1) % 450 == 0:
        batch.commit()
        batch = db.batch()

batch.commit()

# Update metadata
total = len(list(db.collection("products").stream()))
all_collections = set()
all_categories = set()
for doc in db.collection("products").stream():
    d = doc.to_dict()
    if d.get("collection"):
        all_collections.add(d["collection"])
    all_categories.add(d["category"])

db.collection("catalog-meta").document("visconti").set({
    "name": "Visconti Furniture",
    "vendor": "Visconti",
    "totalProducts": total,
    "lastUpdated": now,
    "collections": sorted(all_collections),
    "categories": sorted(all_categories),
})

print(f"\nDone! Total products in Firestore: {total}")
print(f"Collections: {sorted(all_collections)}")
print(f"Categories: {sorted(all_categories)}")
